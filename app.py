from flask import Flask, request, send_file, render_template
import os
import zipfile
from pptx import Presentation
from pptx.util import Inches
import io
import json
import gc

app = Flask(__name__, template_folder='templates')

TEMPLATE_PATH = os.path.join('templates', 'default_template.pptx')

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    uploaded_files = request.files.getlist('images')
    image_map = {f.filename: f.read() for f in uploaded_files}

    project_data = json.loads(request.form.get('project_data', '[]'))

    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w') as zipf:
        for project in project_data:
            *img_filenames, title = project
            prs = Presentation(TEMPLATE_PATH)
            title_text = f"{title} 광고 상품 소개서"

            for shape in prs.slides[0].shapes:
                if shape.has_text_frame and "광고 상품 소개서" in shape.text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "광고 상품 소개서" in run.text:
                                run.text = title_text

            for img_name in img_filenames:
                img_bytes = image_map.get(img_name)
                if not img_bytes:
                    continue
                img_stream = io.BytesIO(img_bytes)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(
                    img_stream, Inches(0), Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )

            ppt_bytes = io.BytesIO()
            prs.save(ppt_bytes)
            ppt_bytes.seek(0)
            zipf.writestr(f"{title_text}.pptx", ppt_bytes.read())

            del prs, ppt_bytes
            gc.collect()

    zip_buffer.seek(0)
    return send_file(
        zip_buffer,
        as_attachment=True,
        download_name='광고소개서_모음.zip',
        mimetype='application/zip'
    )

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)
